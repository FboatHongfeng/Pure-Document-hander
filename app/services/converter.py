"""文件格式转换引擎

扩展接口: ConverterInterface，所有转换器需实现 convert(input, output) 方法
"""
import os
import subprocess
import tempfile
import threading
from abc import ABC, abstractmethod
from pathlib import Path

from app.utils.file_utils import get_extension, safe_filename
from app.services.dependency import find_libreoffice
from app.utils.logger import get_logger

logger = get_logger("converter")

# Windows: 隐藏终端窗口
_HIDE_TERMINAL = 0x08000000 if os.name == "nt" else 0


# ---------- 中文字体注册 ----------

_CJK_FONT_NAME = None

def _get_cjk_font() -> str | None:
    """查找并注册中文字体，返回字体名"""
    global _CJK_FONT_NAME
    if _CJK_FONT_NAME:
        return _CJK_FONT_NAME

    font_paths = [
        ("C:/Windows/Fonts/msyh.ttc", 0),   # 微软雅黑
        ("C:/Windows/Fonts/simsun.ttc", 0),  # 宋体
        ("C:/Windows/Fonts/simhei.ttf", 0),  # 黑体
        ("C:/Windows/Fonts/simkai.ttf", 0),  # 楷体
    ]

    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    for fpath, subfont in font_paths:
        if os.path.exists(fpath):
            try:
                name = Path(fpath).stem
                pdfmetrics.registerFont(TTFont(name, fpath, subfontIndex=subfont))
                _CJK_FONT_NAME = name
                logger.info(f"注册中文字体: {fpath} → {name}")
                return name
            except Exception as e:
                logger.warning(f"注册字体失败 {fpath}: {e}")

    logger.warning("未找到中文字体，中文将显示为空白")
    return None


# ---------- 抽象接口 ----------

class ConverterInterface(ABC):
    """转换器基类"""

    @abstractmethod
    def convert(self, input_path: str, output_path: str, **kwargs) -> bool:
        ...

    @abstractmethod
    def supported_inputs(self) -> list[str]:
        ...

    @abstractmethod
    def supported_outputs(self) -> list[str]:
        ...

    @property
    def name(self) -> str:
        return self.__class__.__name__


# ---------- LibreOffice 通用转换 ----------

class LibreOfficeConverter:

    @staticmethod
    def convert(input_path: str, output_dir: str, output_format: str,
                cancel_event: threading.Event | None = None) -> str | None:
        exe = find_libreoffice()
        if not exe:
            return None
        try:
            proc = subprocess.Popen(
                [exe, "--headless", "--convert-to", output_format,
                 "--outdir", output_dir, input_path],
                stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
                creationflags=_HIDE_TERMINAL,
            )
            while proc.poll() is None:
                if cancel_event and cancel_event.is_set():
                    proc.kill()
                    proc.wait()
                    return None
                try:
                    proc.wait(timeout=0.5)
                except subprocess.TimeoutExpired:
                    pass
            stem = Path(input_path).stem
            for cand in [f"{stem}.{output_format}", f"{stem}.{output_format.lower()}"]:
                out = os.path.join(output_dir, cand)
                if os.path.exists(out):
                    return out
            return None
        except Exception as e:
            logger.warning(f"LibreOffice转换失败: {e}")
            return None


# ---------- reportlab PDF 辅助 ----------

def _create_cjk_canvas(output_path: str, pagesize=None):
    """创建带中文字体的reportlab Canvas"""
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas
    c = canvas.Canvas(output_path, pagesize=pagesize or A4)
    font_name = _get_cjk_font()
    if font_name:
        c.setFont(font_name, 11)
    return c, font_name


# ---------- 具体转换器 ----------

class DocxToPdfConverter(ConverterInterface):
    """Word -> PDF — COM > docx2pdf > LibreOffice > Python"""

    def supported_inputs(self) -> list[str]:
        return [".docx", ".doc"]

    def supported_outputs(self) -> list[str]:
        return [".pdf"]

    def convert(self, input_path: str, output_path: str, **kwargs) -> bool:
        # 1) COM (Word) — 最高保真度
        if self._com_convert(input_path, output_path):
            return True

        # 2) docx2pdf
        try:
            from docx2pdf import convert as docx2pdf_convert
            docx2pdf_convert(input_path, output_path)
            if os.path.exists(output_path):
                return True
        except Exception:
            pass

        # 3) LibreOffice
        lo = LibreOfficeConverter()
        lo_result = lo.convert(input_path, os.path.dirname(output_path), "pdf",
                               cancel_event=kwargs.get("cancel_event"))
        if lo_result:
            if lo_result != output_path:
                import shutil
                shutil.move(lo_result, output_path)
            return True

        # 4) python-docx + reportlab（带图片/表格）
        try:
            self._python_convert(input_path, output_path)
            return os.path.exists(output_path)
        except Exception as e:
            logger.error(f"Word->PDF转换失败: {e}")
            return False

    def _com_convert(self, input_path: str, output_path: str) -> bool:
        """通过 Word COM 实现完整保真转换"""
        try:
            import pythoncom
            import win32com.client
            abs_input = os.path.abspath(input_path)
            abs_output = os.path.abspath(output_path)
            pythoncom.CoInitialize()
            try:
                word = win32com.client.Dispatch("Word.Application")
                try:
                    word.Visible = False
                except Exception:
                    pass
                doc = word.Documents.Open(abs_input, ReadOnly=True)
                doc.SaveAs(abs_output, FileFormat=17)  # 17 = pdf
                doc.Close()
                word.Quit()
            finally:
                pythoncom.CoUninitialize()
            if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                logger.info(f"COM Word->PDF 完成: {output_path}")
                return True
            return False
        except ImportError:
            return False
        except Exception as e:
            logger.warning(f"COM Word转换失败: {e}")
            return False

    def _python_convert(self, input_path: str, output_path: str) -> None:
        from docx import Document
        from reportlab.lib.pagesizes import A4

        doc = Document(input_path)
        c, font_name = _create_cjk_canvas(output_path)
        page_w, page_h = A4
        mx = 55
        mw = page_w - 110        # 可用文字宽度
        y = page_h - 50
        fn = font_name or "Helvetica"

        def wrapped_text(text, fs, leading_extra=0):
            """折行绘制: 用 stringWidth 精确测量中文宽度"""
            nonlocal y
            if y < 50:
                c.showPage()
                y = page_h - 50

            c.setFont(fn, fs)
            line_h = fs * 1.6 + leading_extra
            # 逐字测量, 超出宽度则换行
            line = ""
            for ch in text:
                test = line + ch
                try:
                    sw = c.stringWidth(test, fn, fs)
                except Exception:
                    cjk = sum(1 for ch in test if '一' <= ch <= '鿿')
                    sw = cjk * fs + (len(test) - cjk) * fs * 0.55
                if sw > mw and line:
                    c.drawString(mx, y, line)
                    y -= line_h
                    if y < 50:
                        c.showPage()
                        y = page_h - 50
                        c.setFont(fn, fs)
                    line = ch
                else:
                    line = test
            if line:
                c.drawString(mx, y, line)
                y -= line_h

        # 提取图片
        import zipfile, tempfile
        image_map = {}
        try:
            from xml.etree import ElementTree as ET
            with zipfile.ZipFile(input_path) as zf:
                rels_xml = zf.read('word/_rels/document.xml.rels')
                ns = '{http://schemas.openxmlformats.org/package/2006/relationships}'
                for rel in ET.fromstring(rels_xml):
                    target = rel.get('Target', '')
                    if 'image' in target.lower():
                        img_path = 'word/' + target.replace('../', '')
                        try:
                            image_map[rel.get('Id')] = zf.read(img_path)
                        except KeyError:
                            pass
        except Exception:
            pass

        for para in doc.paragraphs:
            # 检查段落中是否有图片
            has_img = False
            for run in para.runs:
                for rId, img_data in image_map.items():
                    if rId in run._element.xml:
                        try:
                            tmp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
                            tmp.write(img_data); tmp.close()
                            from PIL import Image as PILImg
                            pi = PILImg.open(tmp.name); iw, ih = pi.size
                            ratio = mw / max(iw, 1)
                            rw, rh = iw * ratio, ih * ratio
                            if y < rh + 20:
                                c.showPage()
                                y = page_h - 50
                            c.drawImage(tmp.name, mx, y - rh, rw, rh)
                            y -= rh + 10; os.unlink(tmp.name)
                            has_img = True; break
                        except Exception:
                            pass
                if has_img:
                    break
            if has_img:
                continue

            if not para.text.strip():
                y -= 6
                continue
            is_heading = para.style.name.startswith("Heading") if para.style else False
            fs = 14 if is_heading else 11
            if is_heading:
                y -= 6
            wrapped_text(para.text, fs)
            if is_heading:
                y -= 4

        # 表格
        for table in doc.tables:
            y -= 10
            if y < 50:
                c.showPage()
                y = page_h - 50
            for row in table.rows:
                cells = [cell.text.strip()[:50] for cell in row.cells]
                wrapped_text(" | ".join(cells), 8)
            y -= 6

        c.save()
        logger.info(f"Python Word->PDF 完成: {output_path}")


class PdfToDocxConverter(ConverterInterface):
    """PDF -> Word — Word COM > pdf2docx"""

    def supported_inputs(self) -> list[str]:
        return [".pdf"]

    def supported_outputs(self) -> list[str]:
        return [".docx"]

    def convert(self, input_path: str, output_path: str, **kwargs) -> bool:
        # 1) pdf2docx (对文字型PDF效果最好)
        try:
            from pdf2docx import Converter as Pdf2Docx
            cv = Pdf2Docx(input_path)
            cv.convert(output_path)
            cv.close()
            if os.path.exists(output_path) and os.path.getsize(output_path) > 1024:
                return True
        except Exception as e:
            logger.warning(f"pdf2docx失败: {e}, 尝试COM回退")

        # 2) Word COM (回退方案)
        if self._com_convert(input_path, output_path):
            return True

        logger.error(f"PDF->Docx转换失败: 所有方法均不可用")
        return False

    def _com_convert(self, input_path: str, output_path: str) -> bool:
        try:
            import pythoncom
            import win32com.client
            abs_input = os.path.abspath(input_path)
            abs_output = os.path.abspath(output_path)
            pythoncom.CoInitialize()
            try:
                word = win32com.client.Dispatch("Word.Application")
                try:
                    word.Visible = False
                except Exception:
                    pass
                doc = word.Documents.Open(abs_input, ReadOnly=True)
                doc.SaveAs(abs_output, FileFormat=16)  # 16 = docx
                doc.Close()
                word.Quit()
            finally:
                pythoncom.CoUninitialize()
            if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                logger.info(f"COM PDF->Word 完成: {output_path}")
                return True
            return False
        except ImportError:
            return False
        except Exception as e:
            logger.warning(f"COM PDF->Word失败: {e}")
            return False


class PptxToPdfConverter(ConverterInterface):
    """PPT -> PDF — 优先COM(全保真) > LibreOffice > Python"""

    def supported_inputs(self) -> list[str]:
        return [".pptx", ".ppt"]

    def supported_outputs(self) -> list[str]:
        return [".pdf"]

    def convert(self, input_path: str, output_path: str, **kwargs) -> bool:
        # 方案1: COM (PowerPoint 安装时)
        if self._com_convert(input_path, output_path):
            return True

        # 方案2: LibreOffice
        lo = LibreOfficeConverter()
        lo_result = lo.convert(input_path, os.path.dirname(output_path), "pdf",
                               cancel_event=kwargs.get("cancel_event"))
        if lo_result:
            if lo_result != output_path:
                import shutil
                shutil.move(lo_result, output_path)
            return True

        # 方案3: Python 文字提取
        try:
            self._python_convert(input_path, output_path)
            return os.path.exists(output_path)
        except Exception as e:
            logger.error(f"PPTX->PDF转换失败: {e}")
            return False

    def _com_convert(self, input_path: str, output_path: str) -> bool:
        """通过 PowerPoint COM 实现完整保真转换"""
        try:
            import pythoncom
            import win32com.client
            abs_input = os.path.abspath(input_path)
            abs_output = os.path.abspath(output_path)

            pythoncom.CoInitialize()
            try:
                ppt = win32com.client.Dispatch("PowerPoint.Application")
                try:
                    ppt.Visible = False
                except Exception:
                    pass  # PowerPoint 2016+ 不允许隐藏窗口
                presentation = ppt.Presentations.Open(abs_input, WithWindow=False)
                presentation.SaveAs(abs_output, 32)  # 32 = ppSaveAsPDF
                presentation.Close()
                ppt.Quit()
            finally:
                pythoncom.CoUninitialize()

            if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                logger.info(f"COM PPT->PDF 完成: {output_path}")
                return True
            return False
        except ImportError:
            logger.info("pywin32 未安装，跳过COM转换")
            return False
        except Exception as e:
            logger.warning(f"COM PPT转换失败: {e}")
            return False

    def _python_convert(self, input_path: str, output_path: str) -> None:
        from pptx import Presentation
        from reportlab.lib.pagesizes import A4

        prs = Presentation(input_path)
        c, font_name = _create_cjk_canvas(output_path)
        pw, ph = A4            # ph = page height (只读, 无需nonlocal)
        mx, mw = 55, pw - 110
        fn = font_name or "Helvetica"
        y = ph - 50            # 可变, 在_wrap中用nonlocal

        def _wrap(text, fs):
            nonlocal y
            if y < 45:
                c.showPage()
                y = ph - 50
            c.setFont(fn, fs)
            lh = fs * 1.5
            line = ""
            for ch in text:
                test = line + ch
                try:
                    sw = c.stringWidth(test, fn, fs)
                except Exception:
                    cjk = sum(1 for ch in test if '一' <= ch <= '鿿')
                    sw = cjk * fs + (len(test) - cjk) * fs * 0.55
                if sw > mw and line:
                    c.drawString(mx, y, line)
                    y -= lh
                    if y < 45:
                        c.showPage()
                        y = ph - 50
                        c.setFont(fn, fs)
                    line = ch
                else:
                    line = test
            if line:
                c.drawString(mx, y, line)
                y -= lh

        for slide in prs.slides:
            y = ph - 50
            title_shape = slide.shapes.title
            title_id = title_shape.shape_id if title_shape else None

            if title_shape and title_shape.text:
                _wrap(title_shape.text, 16)
                if y > 45:
                    y -= 8

            for shape in slide.shapes:
                if title_id is not None and shape.shape_id == title_id:
                    continue        # 标题已渲染，跳过
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            _wrap(para.text, 10)
            c.showPage()
        c.save()


class XlsxToPdfConverter(ConverterInterface):
    """Excel -> PDF — COM > LibreOffice > Python"""

    def supported_inputs(self) -> list[str]:
        return [".xlsx", ".xls"]

    def supported_outputs(self) -> list[str]:
        return [".pdf"]

    def convert(self, input_path: str, output_path: str, **kwargs) -> bool:
        # 1) COM (Excel)
        if self._com_convert(input_path, output_path, "pdf"):
            return True
        # 2) LibreOffice
        lo = LibreOfficeConverter()
        lo_result = lo.convert(input_path, os.path.dirname(output_path), "pdf",
                               cancel_event=kwargs.get("cancel_event"))
        if lo_result:
            if lo_result != output_path:
                import shutil
                shutil.move(lo_result, output_path)
            return True
        # 3) Python
        try:
            self._python_convert(input_path, output_path)
            return os.path.exists(output_path)
        except Exception as e:
            logger.error(f"XLSX->PDF转换失败: {e}")
            return False

    def _com_convert(self, input_path, output_path, fmt="pdf"):
        try:
            import pythoncom, win32com.client
            abs_in = os.path.abspath(input_path)
            abs_out = os.path.abspath(output_path)
            pythoncom.CoInitialize()
            try:
                excel = win32com.client.Dispatch("Excel.Application")
                try: excel.Visible = False
                except: pass
                wb = excel.Workbooks.Open(abs_in, ReadOnly=True)
                if fmt == "pdf":
                    wb.ExportAsFixedFormat(0, abs_out)  # 0=pdf
                elif fmt == "docx":
                    wb.SaveAs(abs_out, FileFormat=51)  # 51=xlsx, we use SaveAs for copy
                wb.Close()
                excel.Quit()
            finally:
                pythoncom.CoUninitialize()
            return os.path.exists(abs_out) and os.path.getsize(abs_out) > 0
        except Exception as e:
            logger.warning(f"COM Excel转换失败: {e}")
            return False

    def _python_convert(self, input_path: str, output_path: str) -> None:
        from openpyxl import load_workbook
        from reportlab.lib.pagesizes import A4, landscape

        wb = load_workbook(input_path, data_only=True)
        c, font_name = _create_cjk_canvas(output_path, landscape(A4))
        page_w, page_h = landscape(A4)
        mx, mw = 40, page_w - 80
        fn = font_name or "Helvetica"

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            y = page_h - 40
            c.setFont(fn, 12)
            c.drawString(mx, y, f"Sheet: {sheet_name}")
            y -= 22
            for row in ws.iter_rows(max_row=100, values_only=True):
                row_text = " | ".join(str(cell)[:40] if cell else "" for cell in row[:10])
                c.setFont(fn, 7)
                lh = 10
                line = ""
                for ch in row_text:
                    test = line + ch
                    try:
                        sw = c.stringWidth(test, fn, 7)
                    except Exception:
                        cjk_n = sum(1 for c2 in test if '一' <= c2 <= '鿿')
                        sw = cjk_n * 7 + (len(test) - cjk_n) * 4
                    if sw > mw and line:
                        c.drawString(mx, y, line); y -= lh
                        if y < 35: c.showPage(); y = page_h - 35; c.setFont(fn, 7)
                        line = ch
                    else:
                        line = test
                if line:
                    c.drawString(mx, y, line); y -= lh
                if y < 35:
                    c.showPage()
                    y = page_h - 35
            c.showPage()
        c.save()


class XlsxToDocxConverter(ConverterInterface):
    """Excel -> Word — COM > python-docx"""

    def supported_inputs(self) -> list[str]:
        return [".xlsx", ".xls"]

    def supported_outputs(self) -> list[str]:
        return [".docx"]

    def convert(self, input_path: str, output_path: str, **kwargs) -> bool:
        # 1) COM (Excel SaveAs docx — limited, Excel doesn't directly save to docx)
        # Actually Excel COM can't save to docx. Use Excel→PDF→Word or python-docx
        abs_in = os.path.abspath(input_path)
        abs_out = os.path.abspath(output_path)

        # 2) python-docx: build a table-based Word document from Excel
        try:
            self._python_convert(input_path, output_path)
            return os.path.exists(output_path) and os.path.getsize(output_path) > 0
        except Exception as e:
            logger.error(f"XLSX->DOCX转换失败: {e}")
            return False

    def _python_convert(self, input_path, output_path):
        from openpyxl import load_workbook
        from docx import Document
        from docx.shared import Pt, Inches, Cm
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        wb = load_workbook(input_path, data_only=True)
        doc = Document()

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            doc.add_heading(sheet_name, level=1)
            rows = list(ws.iter_rows(max_row=200, values_only=True))
            if not rows:
                continue

            max_cols = max(len([c for c in r if c is not None]) for r in rows) if rows else 1
            max_cols = min(max_cols, 15)
            table = doc.add_table(rows=len(rows), cols=max_cols, style='Table Grid')

            for i, row in enumerate(rows):
                for j in range(max_cols):
                    val = str(row[j]) if j < len(row) and row[j] is not None else ""
                    cell = table.cell(i, j)
                    cell.text = val[:200]
                    for p in cell.paragraphs:
                        for run in p.runs:
                            run.font.size = Pt(8)

            doc.add_paragraph()  # spacer

        doc.save(output_path)
        logger.info(f"Python XLSX->DOCX 完成: {output_path}")


class ImagesToPdfConverter(ConverterInterface):
    """图片 -> PDF"""

    def supported_inputs(self) -> list[str]:
        return [".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tiff", ".webp"]

    def supported_outputs(self) -> list[str]:
        return [".pdf"]

    def convert(self, input_path: str, output_path: str, **kwargs) -> bool:
        try:
            from PIL import Image
            img = Image.open(input_path)
            if img.mode in ("RGBA", "P"):
                img = img.convert("RGB")
            img.save(output_path, "PDF")
            return True
        except Exception as e:
            logger.error(f"图片->PDF转换失败: {e}")
            return False


# ---------- 转换器注册表 ----------

_converter_registry: dict[str, dict[str, ConverterInterface]] = {}

def _register(converter: ConverterInterface) -> None:
    for in_ext in converter.supported_inputs():
        if in_ext not in _converter_registry:
            _converter_registry[in_ext] = {}
        for out_ext in converter.supported_outputs():
            if in_ext == out_ext:
                continue       # 跳过同类型互转
            _converter_registry[in_ext][out_ext] = converter


_register(DocxToPdfConverter())
_register(PdfToDocxConverter())
_register(PptxToPdfConverter())
_register(XlsxToPdfConverter())
_register(XlsxToDocxConverter())
_register(ImagesToPdfConverter())


# ---------- 公共API ----------

def get_supported_conversions() -> list[dict]:
    result = []
    for in_ext, outputs in _converter_registry.items():
        for out_ext in outputs:
            result.append({
                "input": in_ext,
                "output": out_ext,
                "label": f"{in_ext} -> {out_ext}",
            })
    return result


def get_target_formats(input_ext: str) -> list[str]:
    """给定输入格式，返回可转换的目标格式列表（已排除同类转换）"""
    converters = _converter_registry.get(input_ext, {})
    return [out for out in converters if out != input_ext]


def convert_file(input_path: str, output_path: str, **kwargs) -> tuple[bool, str]:
    in_ext = get_extension(input_path)
    out_ext = get_extension(output_path)

    converters = _converter_registry.get(in_ext, {})
    converter = converters.get(out_ext)

    if not converter:
        return False, f"不支持转换: {in_ext} -> {out_ext}"

    logger.info(f"开始转换: {input_path} -> {output_path}")

    try:
        ok = converter.convert(input_path, output_path, **kwargs)
        if ok:
            logger.info(f"转换成功: {output_path}")
            return True, ""
        logger.error(f"转换失败: {input_path} -> {output_path}")
        return False, "转换失败"
    except Exception as e:
        logger.exception(f"转换异常: {input_path}")
        return False, str(e)


def is_conversion_supported(input_ext: str, output_ext: str) -> bool:
    if input_ext == output_ext:
        return False
    converters = _converter_registry.get(input_ext, {})
    return output_ext in converters


def is_complex_file(input_path: str) -> bool:
    ext = get_extension(input_path)
    try:
        size_mb = os.path.getsize(input_path) / (1024 * 1024)
    except OSError:
        return False

    if ext in (".docx", ".doc"):
        if size_mb > 10:
            return True
        try:
            from docx import Document
            doc = Document(input_path)
            if len(doc.tables) > 5 or size_mb > 5:
                return True
        except Exception:
            pass

    if ext in (".pptx", ".ppt"):
        if size_mb > 20:
            return True
        try:
            from pptx import Presentation
            prs = Presentation(input_path)
            if len(prs.slides) > 30 or size_mb > 10:
                return True
        except Exception:
            pass

    if ext in (".xlsx", ".xls") and size_mb > 5:
        return True

    return False
